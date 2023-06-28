from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor


class PowerPointPresentation:
	def __init__(self, file):
		self.presentation = Presentation()
		self.dict_about_team = file
		self.path_on_disc = self.dict_about_team['path_on_disc']
		self.name_1 = self.dict_about_team['name_1']
		self.name_2 = self.dict_about_team['name_2']
		self.coefficient_1 = self.dict_about_team['coefficient_dict'][self.name_1]
		self.coefficient_2 = self.dict_about_team['coefficient_dict'][self.name_2]
	
	def create_presentation(self):
		self.add_first_slide()
		# save presentation
		self.save_pptx_file()
	
	def add_first_slide(self):
		slide, title_placeholder = self.create_new_slide_and_placeholder()
		# add teams logo
		slide.shapes.add_picture(
			f"{self.path_on_disc}/{self.name_1}.jpg", left=Inches(0.5), top=Inches(2.5), height=Inches(1.3), width=Inches(1.3))
		slide.shapes.add_picture(
			f"{self.path_on_disc}/{self.name_2}.jpg", left=Inches(8.4), top=Inches(2.5), height=Inches(1.3), width=Inches(1.3))
		content_placeholder = slide.placeholders[1]
		title_placeholder.text = f"{self.name_1}  против  {self.name_2}"
		content_placeholder.text = \
			f"Коефициент\n{self.coefficient_1} - {self.coefficient_2}\nКоличество игр {self.dict_about_team['best_of_number']}"
		self.change_title_color(title_placeholder)
		self.change_content_color(content_placeholder)
		self.change_background_color(slide)
		
	def add_second_slide(self):
		pass
	
	def create_new_slide_and_placeholder(self):
		# create slide
		slide_layout = self.presentation.slide_layouts[0]
		slide = self.presentation.slides.add_slide(slide_layout)
		title_placeholder = slide.shapes.title
		return slide, title_placeholder
	
	@staticmethod
	def change_title_color(title_placeholder):
		# title color
		font = title_placeholder.text_frame.paragraphs[0].runs[0].font
		font.color.rgb = RGBColor(255, 255, 255)
	
	@staticmethod
	def change_content_color(content_placeholder):
		# content color
		paragraphs_content = content_placeholder.text_frame.paragraphs
		for paragraph in paragraphs_content:
			font = paragraph.runs[0].font
			font.color.rgb = RGBColor(255, 255, 255)
	
	@staticmethod
	def change_background_color(slide):
		# background color
		slide.background.fill.solid()
		slide.background.fill.fore_color.rgb = RGBColor(25, 43, 69)
		
	def save_pptx_file(self):
		self.presentation.save(f'{self.path_on_disc}/{self.name_1}_vs_{self.name_2}.pptx')
